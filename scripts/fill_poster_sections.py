from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

base_dir = Path(r"c:\Users\Youssef\OneDrive\Desktop\esprithub")
input_pptx = base_dir / "template weita_POSTERSExemple.pptx"
output_pptx = base_dir / "EspriHub_Poster_Filled.pptx"
output_pdf = base_dir / "EspriHub_Poster_Filled.pdf"

approach_text = (
    "System design and workflow:\n"
    "1) Full-stack architecture: Angular frontend + Spring Boot backend + PostgreSQL/Redis via Docker for scalable classroom deployment.\n"
    "2) Academic domain model: role-based modules (ADMIN, CHIEF, TEACHER, STUDENT), task entities, group/class assignment, and submission lifecycle tracking.\n"
    "3) GitHub integration pipeline: OAuth authentication, repository linking, webhook subscription, signature validation (HMAC SHA-256), and event-driven data sync.\n"
    "4) Automation layer: schedulers for webhook health checks, failed-subscription retry, and stale webhook cleanup to ensure reliability over time.\n"
    "5) AI-assisted evaluation support: submission file ingestion, commit-context visibility, and teacher-centered grading support (human-in-the-loop).\n"
    "6) Governance and security: JWT stateless auth, endpoint authorization by role, secure headers, and validation of OAuth state and institutional email constraints."
)

results_text = (
    "Evidence from implemented prototype and platform behavior:\n"
    "1) End-to-end operational prototype delivered: repository browsing, commit/history tracking, submission workflows, notifications, and admin management are integrated in one platform.\n"
    "2) Webhook reliability improved through self-healing operations: automatic subscription, periodic health checks, retry policies, and cleanup routines reduce manual intervention needs.\n"
    "3) Monitoring visibility strengthened: admin views expose webhook status, failures, and repository-level synchronization signals for faster troubleshooting.\n"
    "4) Assessment traceability increased: task assignment, submission attempts, commit hashes, and grading metadata are centrally recorded for auditability.\n"
    "5) Risk analytics capability added: a Deadline Risk Radar endpoint aggregates overdue tasks, due-soon tasks, repository staleness, token availability, and inactivity signals.\n"
    "6) Security posture established for academic deployment: signature-verified webhooks, role-based access control, JWT flows, and production hardening checklist documented.\n"
    "7) Practical impact observed (qualitative): reduced repetitive manual coordination for teachers, improved transparency of student contribution signals, and better readiness for large-cohort project monitoring."
)

prs = Presentation(str(input_pptx))
slide = prs.slides[0]

# Shape IDs discovered via inspection script.
APPROACH_SHAPE_INDEX = 21  # no-text placeholder under "Approach & Methodology"
RESULTS_SHAPE_INDEX = 42   # no-text placeholder under "Results & Evidence"

approach_shape = slide.shapes[APPROACH_SHAPE_INDEX - 1]
results_shape = slide.shapes[RESULTS_SHAPE_INDEX - 1]

for shape, text in ((approach_shape, approach_text), (results_shape, results_text)):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    # Keep typography compact for poster boxes.
    for run in p.runs:
        run.font.size = Pt(14)

prs.save(str(output_pptx))

# Export to PDF via PowerPoint COM if available.
pdf_exported = False
try:
    import comtypes.client  # type: ignore

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(str(output_pptx), WithWindow=False)
    # 32 = ppSaveAsPDF
    presentation.SaveAs(str(output_pdf), 32)
    presentation.Close()
    powerpoint.Quit()
    pdf_exported = True
except Exception:
    pdf_exported = False

print(f"PPTX written: {output_pptx}")
print(f"PDF written: {output_pdf}" if pdf_exported else "PDF export skipped (PowerPoint COM unavailable)")
