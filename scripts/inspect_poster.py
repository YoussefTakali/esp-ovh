from pathlib import Path
from pptx import Presentation

pptx_path = Path(r"c:\Users\Youssef\OneDrive\Desktop\esprithub\template weita_POSTERSExemple.pptx")
prs = Presentation(str(pptx_path))

print(f"Slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides, start=1):
    print(f"\n--- Slide {i} ---")
    for j, shape in enumerate(slide.shapes, start=1):
        left = getattr(shape, "left", None)
        top = getattr(shape, "top", None)
        width = getattr(shape, "width", None)
        height = getattr(shape, "height", None)
        txt = ""
        if hasattr(shape, "text"):
            txt = (shape.text or "").strip().replace("\n", " | ")
        if txt:
            print(f"{j:02d} pos=({left},{top},{width},{height}) text={txt[:260]}")
        else:
            print(f"{j:02d} pos=({left},{top},{width},{height}) [no text]")
